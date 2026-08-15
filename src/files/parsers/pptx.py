"""Bounded parser for Microsoft PowerPoint (.pptx)."""

from src.core.config import settings


class PPTXParser:
    @staticmethod
    def parse_pptx(file_path: str) -> tuple[str | None, bool]:
        try:
            from pptx import Presentation

            presentation = Presentation(file_path)
            parts: list[str] = []
            total_chars = 0
            limit = settings.max_document_extract_chars
            for idx, slide in enumerate(presentation.slides, start=1):
                if total_chars >= limit:
                    break
                slide_parts = [f"--- Slide {idx} ---"]
                for shape in slide.shapes:
                    text = getattr(shape, "text", "")
                    if text and text.strip():
                        slide_parts.append(text.strip())
                block = "\n".join(slide_parts)
                remaining = max(0, limit - total_chars)
                parts.append(block[:remaining])
                total_chars += min(len(block), remaining)

            full_text = "\n\n".join(parts).strip()
            return full_text if full_text else None, True
        except Exception as exc:
            return f"Error parsing PPTX: {exc}", False


pptx_parser = PPTXParser()
